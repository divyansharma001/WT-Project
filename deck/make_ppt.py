"""
Generate haven.pptx — a 16:9 deck matching the landing-page theme:
serene sky-blue palette, Poppins typography, rounded cards, soft fills.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ------------------------------------------------------------------
# Palette (mirrors --serene-*, --slate-*, --amber-* from style.css)
# ------------------------------------------------------------------
SERENE_50  = RGBColor(0xf0, 0xf9, 0xff)
SERENE_100 = RGBColor(0xe0, 0xf2, 0xfe)
SERENE_200 = RGBColor(0xba, 0xe6, 0xfd)
SERENE_300 = RGBColor(0x7d, 0xd3, 0xfc)
SERENE_400 = RGBColor(0x38, 0xbd, 0xf8)
SERENE_500 = RGBColor(0x0e, 0xa5, 0xe9)
SERENE_600 = RGBColor(0x02, 0x84, 0xc7)
SERENE_700 = RGBColor(0x03, 0x69, 0xa1)
SERENE_800 = RGBColor(0x07, 0x59, 0x85)

SLATE_50  = RGBColor(0xf8, 0xfa, 0xfc)
SLATE_100 = RGBColor(0xf1, 0xf5, 0xf9)
SLATE_200 = RGBColor(0xe2, 0xe8, 0xf0)
SLATE_400 = RGBColor(0x94, 0xa3, 0xb8)
SLATE_500 = RGBColor(0x64, 0x74, 0x8b)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_800 = RGBColor(0x1e, 0x29, 0x3b)

AMBER_50  = RGBColor(0xff, 0xfb, 0xeb)
AMBER_100 = RGBColor(0xfe, 0xf3, 0xc7)
AMBER_200 = RGBColor(0xfd, 0xe6, 0x8a)
AMBER_600 = RGBColor(0xd9, 0x77, 0x06)
AMBER_800 = RGBColor(0x92, 0x40, 0x0e)

PINK_400   = RGBColor(0xf9, 0xa8, 0xd4)
PINK_600   = RGBColor(0xec, 0x48, 0x99)
GREEN_400  = RGBColor(0x86, 0xef, 0xac)
GREEN_600  = RGBColor(0x16, 0xa3, 0x4a)
PURPLE_400 = RGBColor(0xc4, 0xb5, 0xfd)
PURPLE_600 = RGBColor(0x7c, 0x3a, 0xed)

WHITE = RGBColor(0xff, 0xff, 0xff)

FONT = "Poppins"

# ------------------------------------------------------------------
# Presentation setup (16:9)
# ------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def _no_shadow(shape):
    """Remove inherited theme shadow from an autoshape."""
    sp = shape.fill._xPr  # spPr element
    # Remove any existing effectLst, then add empty one
    for el in sp.findall(qn("a:effectLst")):
        sp.remove(el)
    sp.append(etree.SubElement(sp, qn("a:effectLst")))

def bg(slide, color):
    """Fill the slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    _no_shadow(s)
    return s

def rrect(slide, x, y, w, h, fill_color, line_color=None, radius=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.75)
    _no_shadow(s)
    return s

def oval(slide, x, y, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    _no_shadow(s)
    return s

def text(slide, x, y, w, h, content, *,
         size=18, bold=False, color=SLATE_800, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, italic=False, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # Support multi-line input by splitting on '\n'
    lines = content.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def rich_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """runs = list of (text, dict of style overrides) or list of paragraphs where each paragraph is a list of runs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    # runs may be list of (str, style) or list of lists (paragraphs)
    if runs and isinstance(runs[0], list):
        paragraphs = runs
    else:
        paragraphs = [runs]
    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for run_text, style in para:
            r = p.add_run()
            r.text = run_text
            r.font.name = style.get("font", FONT)
            r.font.size = Pt(style.get("size", 16))
            r.font.bold = style.get("bold", False)
            r.font.italic = style.get("italic", False)
            r.font.color.rgb = style.get("color", SLATE_700)
    return tb

def logo(slide, x, y, size=Inches(0.9)):
    """Teardrop-ish logo made from a circle + smaller circle highlight + dot."""
    # Large drop (serene-500)
    big = oval(slide, x, y + size * 0.1, size * 0.9, size * 0.9, SERENE_500)
    # Highlight oval overlapping (serene-400)
    oval(slide, x + size * 0.3, y + size * 0.5, size * 0.6, size * 0.4, SERENE_400)
    # Small dot at top (serene-300)
    dot_size = size * 0.18
    oval(slide, x + size * 0.4, y, dot_size, dot_size, SERENE_300)
    return big

def page_number(slide, n, total):
    text(slide, Inches(12.4), Inches(7.08), Inches(0.9), Inches(0.3),
         f"{n} / {total}", size=9, color=SLATE_400, align=PP_ALIGN.RIGHT)

def brand_footer(slide):
    """Tiny brand mark bottom-left on non-cover slides."""
    logo(slide, Inches(0.55), Inches(7.05), size=Inches(0.28))
    text(slide, Inches(0.92), Inches(7.05), Inches(2.5), Inches(0.3),
         "haven", size=11, bold=True, color=SLATE_700)
    text(slide, Inches(1.55), Inches(7.08), Inches(3), Inches(0.28),
         "  \u00b7  a gentle space for your mind", size=9, color=SLATE_400)

def eyebrow(slide, x, y, label, color=SERENE_600):
    text(slide, x, y, Inches(3), Inches(0.3),
         label.upper(), size=10, bold=True, color=color)

def title(slide, x, y, w, content, color=SLATE_800, size=40):
    return text(slide, x, y, w, Inches(1.2), content,
                size=size, bold=True, color=color, line_spacing=1.1)

def subtitle(slide, x, y, w, content, color=SLATE_500, size=16):
    return text(slide, x, y, w, Inches(1.5), content,
                size=size, color=color, line_spacing=1.45)

# ------------------------------------------------------------------
# Slide builders
# ------------------------------------------------------------------

SLIDES = []

def slide(fn):
    SLIDES.append(fn)
    return fn


@slide
def cover(s, n, total):
    # Soft background
    bg(s, SERENE_50)
    # Decorative blob on right
    oval(s, Inches(9.5), Inches(-2), Inches(7), Inches(7), SERENE_100)
    oval(s, Inches(11), Inches(2), Inches(5), Inches(5), SERENE_200)
    # Logo + wordmark
    logo(s, Inches(0.9), Inches(0.95), size=Inches(0.8))
    text(s, Inches(1.85), Inches(1.0), Inches(4), Inches(0.6),
         "haven", size=28, bold=True, color=SLATE_800)
    # Hero copy
    rich_text(s, Inches(0.9), Inches(2.4), Inches(9), Inches(3.5),
              [
                [("A gentle space\nfor your ",
                  {"size": 68, "bold": True, "color": SLATE_800})],
                [("mind.", {"size": 68, "bold": True, "color": SERENE_600})],
              ],
              line_spacing=1.05)
    text(s, Inches(0.9), Inches(5.6), Inches(8), Inches(0.6),
         "An AI-powered mental wellness companion.",
         size=18, color=SLATE_500)
    # Bottom pill
    rrect(s, Inches(0.9), Inches(6.4), Inches(3.1), Inches(0.45),
          SERENE_100, radius=0.5)
    oval(s, Inches(1.05), Inches(6.52), Inches(0.22), Inches(0.22), SERENE_500)
    text(s, Inches(1.4), Inches(6.48), Inches(3), Inches(0.32),
         "Project presentation  \u00b7  2026",
         size=11, bold=True, color=SERENE_800)
    # No brand footer on cover
    page_number(s, n, total)


@slide
def problem(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "The problem")
    title(s, Inches(0.9), Inches(1.15), Inches(11), "Mental health help is\nhard to reach when you need it most.")
    # Stat cards
    cards = [
        ("1 in 5", "adults experience mental illness in any given year.", SERENE_600),
        ("$150+", "average cost per therapy session \u2014 often out of pocket.", PINK_600),
        ("2 a.m.", "is when a lot of spirals happen. Therapists are asleep.", PURPLE_600),
        ("80%+", "of mental-health apps sell or share user data.", AMBER_600),
    ]
    x = Inches(0.9)
    y = Inches(3.4)
    w = Inches(2.85)
    h = Inches(2.7)
    gap = Inches(0.2)
    for i, (stat, desc, color) in enumerate(cards):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, WHITE, line_color=SLATE_100, radius=0.08)
        text(s, cx + Inches(0.3), y + Inches(0.35), w - Inches(0.6), Inches(0.9),
             stat, size=34, bold=True, color=color)
        text(s, cx + Inches(0.3), y + Inches(1.4), w - Inches(0.6), Inches(1.2),
             desc, size=13, color=SLATE_500, line_spacing=1.45)
    brand_footer(s)
    page_number(s, n, total)


@slide
def solution(s, n, total):
    bg(s, SERENE_50)
    oval(s, Inches(-2), Inches(4), Inches(8), Inches(8), SERENE_100)
    eyebrow(s, Inches(0.9), Inches(0.85), "Our solution")
    title(s, Inches(0.9), Inches(1.15), Inches(8.5),
          "Introducing haven.", color=SLATE_800, size=52)
    text(s, Inches(0.9), Inches(2.4), Inches(8.2), Inches(2),
         "A private, browser-based companion that listens, reflects,\nand gently guides \u2014 whenever you need it.",
         size=20, color=SLATE_700, line_spacing=1.5)
    # Three pill tags
    pills = [("Eight mentors", SERENE_600),
             ("Mood-matched", PINK_600),
             ("100% private", GREEN_600)]
    x = Inches(0.9)
    for label, color in pills:
        w = Inches(0.3 + 0.13 * len(label))
        rrect(s, x, Inches(4.7), w, Inches(0.52), WHITE, line_color=SLATE_100, radius=0.5)
        text(s, x + Inches(0.25), Inches(4.78), w, Inches(0.4),
             label, size=13, bold=True, color=color)
        x += w + Inches(0.15)

    # Right: circle with logo
    oval(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(3.6), WHITE)
    oval(s, Inches(9.5), Inches(1.7), Inches(3.0), Inches(3.0), SERENE_100)
    logo(s, Inches(10.3), Inches(2.15), size=Inches(1.4))
    # Floating mentor dots
    for cx, cy, col in [
        (Inches(8.8), Inches(4.5), SERENE_400),
        (Inches(12.6), Inches(3.8), PINK_400),
        (Inches(12.2), Inches(1.3), PURPLE_400),
        (Inches(9.4), Inches(5.4), AMBER_600),
    ]:
        oval(s, cx, cy, Inches(0.5), Inches(0.5), col)
    brand_footer(s)
    page_number(s, n, total)


@slide
def principles(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "What we stand for")
    title(s, Inches(0.9), Inches(1.15), Inches(11),
          "Three non-negotiables.")
    items = [
        ("01", "Privacy is a promise",
         "No server. No database. No analytics. The conversation lives in your browser and leaves with you when you close the tab."),
        ("02", "Boundaries are care",
         "Every mentor knows what it is and what it isn\u2019t. We won\u2019t diagnose or prescribe \u2014 we\u2019ll point to humans when a human is needed."),
        ("03", "Gentleness over cleverness",
         "Mentors are tuned to slow down, reflect, and ask \u2014 not to dispense quick fixes. Soft and useful beats impressive and brittle."),
    ]
    x = Inches(0.9)
    y = Inches(2.9)
    w = Inches(3.85)
    h = Inches(3.7)
    gap = Inches(0.2)
    for i, (num, heading, body) in enumerate(items):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, WHITE, line_color=SLATE_100, radius=0.08)
        text(s, cx + Inches(0.4), y + Inches(0.35), w - Inches(0.8), Inches(0.5),
             num, size=14, bold=True, color=SERENE_600)
        text(s, cx + Inches(0.4), y + Inches(0.95), w - Inches(0.8), Inches(1.0),
             heading, size=22, bold=True, color=SLATE_800, line_spacing=1.2)
        text(s, cx + Inches(0.4), y + Inches(2.1), w - Inches(0.8), Inches(1.5),
             body, size=13, color=SLATE_500, line_spacing=1.55)
    brand_footer(s)
    page_number(s, n, total)


@slide
def mentors(s, n, total):
    bg(s, SERENE_50)
    eyebrow(s, Inches(0.9), Inches(0.55), "Meet your mentors")
    title(s, Inches(0.9), Inches(0.85), Inches(11),
          "Eight voices. One that\u2019s right for now.", size=34)
    subtitle(s, Inches(0.9), Inches(1.85), Inches(11),
             "Each mentor has a distinct style, pace, and specialty. Switch anytime.")
    mentors_list = [
        ("Sarah",  "Compassionate Listener", SERENE_400),
        ("Marcus", "Coping Strategist",      AMBER_600),
        ("Luna",   "Mindfulness Coach",      PURPLE_400),
        ("Kai",    "Motivation & Growth",    GREEN_400),
        ("Elena",  "Relationship Guide",     PINK_400),
        ("David",  "Burnout & Stress",       SERENE_600),
        ("Sofia",  "Self-Compassion",        AMBER_200),
        ("Oliver", "Sleep & Serenity",       SLATE_400),
    ]
    cols = 4
    card_w = Inches(2.85)
    card_h = Inches(1.85)
    gap_x = Inches(0.18)
    gap_y = Inches(0.22)
    start_x = Inches(0.9)
    start_y = Inches(2.75)
    for i, (name, role, color) in enumerate(mentors_list):
        r = i // cols
        c = i % cols
        cx = start_x + (card_w + gap_x) * c
        cy = start_y + (card_h + gap_y) * r
        rrect(s, cx, cy, card_w, card_h, WHITE, line_color=SLATE_100, radius=0.1)
        # Avatar circle
        oval(s, cx + Inches(0.3), cy + Inches(0.3), Inches(1.2), Inches(1.2), color)
        # Initial letter
        text(s, cx + Inches(0.3), cy + Inches(0.53), Inches(1.2), Inches(0.8),
             name[0], size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Name + role
        text(s, cx + Inches(1.65), cy + Inches(0.4), card_w - Inches(1.8), Inches(0.45),
             name, size=16, bold=True, color=SLATE_800)
        text(s, cx + Inches(1.65), cy + Inches(0.92), card_w - Inches(1.8), Inches(0.7),
             role, size=11, color=SERENE_600, line_spacing=1.25)
    brand_footer(s)
    page_number(s, n, total)


@slide
def mood_matching(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "Mood-based matching")
    title(s, Inches(0.9), Inches(1.15), Inches(11),
          "Tell haven how you feel \u2014\nwe\u2019ll suggest the right mentor.")
    rows = [
        ("Happy",   "Kai",    "Motivation & Growth",     GREEN_400),
        ("Calm",    "Luna",   "Mindfulness Coach",       PURPLE_400),
        ("Anxious", "Marcus", "Coping Strategist",       AMBER_600),
        ("Sad",     "Sarah",  "Compassionate Listener",  SERENE_400),
        ("Angry",   "Marcus", "Coping Strategist",       AMBER_600),
        ("Tired",   "Oliver", "Sleep & Serenity",        SLATE_400),
    ]
    x = Inches(0.9)
    y = Inches(3.35)
    w = Inches(5.65)
    h = Inches(0.55)
    gap = Inches(0.12)
    for i, (mood, mentor, role, color) in enumerate(rows):
        col = i // 3
        row = i % 3
        rx = x + (w + Inches(0.35)) * col
        ry = y + (h + gap) * row
        rrect(s, rx, ry, w, h, SLATE_50, line_color=SLATE_100, radius=0.3)
        # Mood tag
        rrect(s, rx + Inches(0.2), ry + Inches(0.1), Inches(1.0), Inches(0.35),
              color, radius=0.5)
        text(s, rx + Inches(0.2), ry + Inches(0.13), Inches(1.0), Inches(0.3),
             mood, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Arrow
        text(s, rx + Inches(1.3), ry + Inches(0.12), Inches(0.4), Inches(0.3),
             "\u2192", size=16, color=SLATE_400)
        # Mentor name
        text(s, rx + Inches(1.7), ry + Inches(0.13), Inches(1.5), Inches(0.3),
             mentor, size=13, bold=True, color=SLATE_800)
        # Role
        text(s, rx + Inches(2.7), ry + Inches(0.15), Inches(3), Inches(0.3),
             role, size=11, color=SLATE_500)
    # Helper tip
    rrect(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.8),
          SERENE_50, line_color=SERENE_100, radius=0.15)
    text(s, Inches(1.2), Inches(5.85), Inches(11), Inches(0.4),
         "You can always switch mentors manually \u2014 the mood pick is a starting point, not a commitment.",
         size=13, color=SERENE_800, italic=True)
    brand_footer(s)
    page_number(s, n, total)


FEATURES = [
    ("Mood-aware matching",
     "Tell haven how you feel and we\u2019ll suggest the right mentor for this moment.",
     SERENE_400, SERENE_700),
    ("Eight distinct mentors",
     "From compassionate listeners to coping strategists \u2014 each with a unique voice.",
     PINK_400, PINK_600),
    ("Evidence-based tools",
     "Gentle CBT reframing, grounding exercises, breathwork, and self-compassion practices.",
     GREEN_400, GREEN_600),
    ("Private by design",
     "No accounts. No database. No tracking. Your key lives only on your device.",
     PURPLE_400, PURPLE_600),
    ("Crisis-aware",
     "Recognizes when more than a chat is needed and surfaces 988 and Crisis Text Line.",
     AMBER_200, AMBER_600),
    ("Export when you want",
     "Download any conversation as a transcript \u2014 revisit, share, or keep your own journal.",
     SLATE_200, SLATE_500),
]


@slide
def features(s, n, total):
    bg(s, SERENE_50)
    eyebrow(s, Inches(0.9), Inches(0.55), "Features")
    title(s, Inches(0.9), Inches(0.85), Inches(11),
          "Built for how you actually feel.", size=34)
    subtitle(s, Inches(0.9), Inches(1.8), Inches(11),
             "Everything designed around one idea: meeting you where you are.")
    cols = 3
    card_w = Inches(3.85)
    card_h = Inches(2.15)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(0.9)
    start_y = Inches(2.75)
    for i, (heading, body, c1, c2) in enumerate(FEATURES):
        r = i // cols
        c = i % cols
        cx = start_x + (card_w + gap_x) * c
        cy = start_y + (card_h + gap_y) * r
        rrect(s, cx, cy, card_w, card_h, WHITE, line_color=SLATE_100, radius=0.08)
        # icon chip
        rrect(s, cx + Inches(0.3), cy + Inches(0.3), Inches(0.55), Inches(0.55),
              c2, radius=0.25)
        # small inner dot to suggest icon
        oval(s, cx + Inches(0.42), cy + Inches(0.42), Inches(0.32), Inches(0.32), c1)
        text(s, cx + Inches(1), cy + Inches(0.35), card_w - Inches(1.2), Inches(0.5),
             heading, size=16, bold=True, color=SLATE_800)
        text(s, cx + Inches(0.3), cy + Inches(1.1), card_w - Inches(0.6), Inches(1.0),
             body, size=12, color=SLATE_500, line_spacing=1.55)
    brand_footer(s)
    page_number(s, n, total)


@slide
def how_it_works(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "How it works")
    title(s, Inches(0.9), Inches(1.15), Inches(11),
          "Three steps to a lighter moment.")
    steps = [
        ("1", "Check in with your mood",
         "Pick from six moods \u2014 or skip it. Your call."),
        ("2", "Choose a mentor",
         "We recommend one based on your mood. Explore all eight freely."),
        ("3", "Start talking",
         "Share what\u2019s on your mind. Your mentor listens, reflects, and guides."),
    ]
    x = Inches(0.9)
    y = Inches(3.3)
    w = Inches(3.85)
    h = Inches(3.2)
    gap = Inches(0.2)
    for i, (num, heading, body) in enumerate(steps):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, WHITE, line_color=SLATE_100, radius=0.08)
        # Big numbered circle
        oval(s, cx + Inches(0.4), y + Inches(0.45), Inches(0.75), Inches(0.75), SERENE_100)
        text(s, cx + Inches(0.4), y + Inches(0.56), Inches(0.75), Inches(0.55),
             num, size=22, bold=True, color=SERENE_700, align=PP_ALIGN.CENTER)
        text(s, cx + Inches(0.4), y + Inches(1.5), w - Inches(0.8), Inches(0.8),
             heading, size=18, bold=True, color=SLATE_800, line_spacing=1.2)
        text(s, cx + Inches(0.4), y + Inches(2.35), w - Inches(0.8), Inches(1.0),
             body, size=13, color=SLATE_500, line_spacing=1.55)
        # Arrow between steps
        if i < 2:
            text(s, cx + w + Inches(0.02), y + Inches(1.3), gap, Inches(0.6),
                 "\u2192", size=28, color=SERENE_300, align=PP_ALIGN.CENTER)
    brand_footer(s)
    page_number(s, n, total)


@slide
def breathing(s, n, total):
    bg(s, SERENE_50)
    # soft circle decoration
    oval(s, Inches(9), Inches(1.5), Inches(5.5), Inches(5.5), SERENE_100)
    oval(s, Inches(9.7), Inches(2.2), Inches(4), Inches(4), SERENE_200)
    oval(s, Inches(10.4), Inches(2.9), Inches(2.6), Inches(2.6), SERENE_500)
    text(s, Inches(10.4), Inches(3.85), Inches(2.6), Inches(0.6),
         "Inhale", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    eyebrow(s, Inches(0.9), Inches(0.85), "Breathe")
    title(s, Inches(0.9), Inches(1.15), Inches(8.2),
          "Four techniques.\nTwo minutes.\nOne steadier you.", size=34)
    subtitle(s, Inches(0.9), Inches(3.35), Inches(8),
             "Interactive guided breathing with real-time visual pacing.")
    # Technique tiles
    techs = [
        ("Box",        "4\u00b74\u00b74\u00b74", "Focus + stress reset",   SERENE_600),
        ("4\u00b77\u00b78", "Relaxing",            "Unwind before sleep",    PURPLE_600),
        ("Coherent",   "5\u00b75",                 "Daily nervous-system reset", GREEN_600),
        ("Extended",   "4\u00b76",                 "Dial down acute anxiety", PINK_600),
    ]
    x = Inches(0.9)
    y = Inches(4.6)
    w = Inches(1.95)
    h = Inches(1.85)
    gap = Inches(0.1)
    for i, (name, pattern, when, color) in enumerate(techs):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, WHITE, line_color=SLATE_100, radius=0.1)
        # color strip
        rrect(s, cx + Inches(0.2), y + Inches(0.25), Inches(0.55), Inches(0.3),
              color, radius=0.4)
        text(s, cx + Inches(0.2), y + Inches(0.26), Inches(0.55), Inches(0.3),
             pattern, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, cx + Inches(0.2), y + Inches(0.75), w - Inches(0.4), Inches(0.5),
             name, size=16, bold=True, color=SLATE_800)
        text(s, cx + Inches(0.2), y + Inches(1.15), w - Inches(0.4), Inches(0.6),
             when, size=10, color=SLATE_500, line_spacing=1.4)
    brand_footer(s)
    page_number(s, n, total)


@slide
def resources(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "Resources")
    title(s, Inches(0.9), Inches(1.15), Inches(11),
          "Tools for the hard moments.")
    subtitle(s, Inches(0.9), Inches(2.2), Inches(11),
             "Curated crisis support, grounding techniques, and trusted reading \u2014 always one tap away.")

    # Crisis banner
    rrect(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(1.8),
          AMBER_50, line_color=AMBER_200, radius=0.08)
    rrect(s, Inches(1.15), Inches(3.4), Inches(0.6), Inches(0.6), AMBER_600, radius=0.2)
    text(s, Inches(1.15), Inches(3.45), Inches(0.6), Inches(0.5),
         "!", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(s, Inches(1.95), Inches(3.4), Inches(9), Inches(0.5),
         "Crisis support, always available", size=18, bold=True, color=AMBER_800)
    text(s, Inches(1.95), Inches(3.9), Inches(10.5), Inches(0.9),
         "988 Suicide & Crisis Lifeline   \u00b7   Text HOME to 741741   \u00b7   SAMHSA 1-800-662-HELP   \u00b7   Trevor Project 1-866-488-7386",
         size=13, color=SLATE_700, line_spacing=1.55)

    # Category summary
    cats = [
        ("Grounding", "5\u00b74\u00b73\u00b72\u00b71 \u00b7 Body scan \u00b7 Anchor object"),
        ("Daily practice", "Three good things \u00b7 Morning light \u00b7 Walk without phone"),
        ("Trusted sources", "NIMH \u00b7 NAMI \u00b7 Mind UK \u00b7 Find A Helpline"),
    ]
    x = Inches(0.9)
    y = Inches(5.2)
    w = Inches(3.85)
    h = Inches(1.4)
    gap = Inches(0.2)
    for i, (head, items) in enumerate(cats):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, SLATE_50, line_color=SLATE_100, radius=0.08)
        text(s, cx + Inches(0.35), y + Inches(0.25), w - Inches(0.7), Inches(0.4),
             head, size=14, bold=True, color=SERENE_700)
        text(s, cx + Inches(0.35), y + Inches(0.7), w - Inches(0.7), Inches(0.8),
             items, size=11, color=SLATE_500, line_spacing=1.55)
    brand_footer(s)
    page_number(s, n, total)


@slide
def privacy(s, n, total):
    # Dark serene background
    bg(s, SERENE_800)
    # soft glow
    oval(s, Inches(9), Inches(-2), Inches(7), Inches(7), SERENE_700)
    oval(s, Inches(11), Inches(1), Inches(5), Inches(5), SERENE_600)
    eyebrow(s, Inches(0.9), Inches(0.85), "Privacy", color=SERENE_200)
    title(s, Inches(0.9), Inches(1.15), Inches(9),
          "Your words. Your device.\nThat\u2019s it.", color=WHITE)
    subtitle(s, Inches(0.9), Inches(3.1), Inches(9),
             "haven runs entirely in your browser. Nothing is stored on any server we control.",
             color=SERENE_100, size=17)

    items = [
        ("No server", "Direct browser \u2192 AI call. No middleman."),
        ("No accounts", "Nothing to sign up for, nothing to be stolen."),
        ("Your key, your rules", "Gemini key stored locally, never transmitted to us."),
        ("Session-only memory", "Close the tab and it\u2019s gone. Truly."),
    ]
    x = Inches(0.9)
    y = Inches(4.55)
    w = Inches(2.85)
    h = Inches(1.9)
    gap = Inches(0.2)
    for i, (head, body) in enumerate(items):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, SERENE_700, radius=0.1)
        # check icon
        oval(s, cx + Inches(0.3), cy_ := y + Inches(0.3), Inches(0.45), Inches(0.45), SERENE_300)
        text(s, cx + Inches(0.3), cy_ + Inches(0.05), Inches(0.45), Inches(0.35),
             "\u2713", size=16, bold=True, color=SERENE_800, align=PP_ALIGN.CENTER)
        text(s, cx + Inches(0.3), y + Inches(0.95), w - Inches(0.6), Inches(0.5),
             head, size=14, bold=True, color=WHITE)
        text(s, cx + Inches(0.3), y + Inches(1.35), w - Inches(0.6), Inches(0.5),
             body, size=11, color=SERENE_100, line_spacing=1.45)
    page_number(s, n, total)


@slide
def architecture(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "Under the hood")
    title(s, Inches(0.9), Inches(1.15), Inches(11),
          "How it works, plainly.")
    subtitle(s, Inches(0.9), Inches(2.2), Inches(11),
             "No black box. The whole data flow is three hops \u2014 and one of them is just rendering.")

    # Flow: Browser -> Gemini -> Browser
    boxes = [
        ("Your browser",   "The whole app (HTML, CSS, JS)\nloads here. No backend.",     SERENE_500),
        ("Gemini API",     "Direct call to Google\u2019s\nGemini 2.0 Flash endpoint.",   PURPLE_600),
        ("Your browser",   "Response renders as a\nchat message. History in memory.",   SERENE_500),
    ]
    x = Inches(0.9)
    y = Inches(3.4)
    w = Inches(3.6)
    h = Inches(2.6)
    gap = Inches(0.45)
    for i, (head, body, color) in enumerate(boxes):
        cx = x + (w + gap) * i
        rrect(s, cx, y, w, h, WHITE, line_color=SLATE_100, radius=0.08)
        rrect(s, cx, y, w, Inches(0.1), color, radius=0.5)
        text(s, cx + Inches(0.35), y + Inches(0.45), w - Inches(0.7), Inches(0.5),
             f"Step {i+1}", size=11, bold=True, color=color)
        text(s, cx + Inches(0.35), y + Inches(0.95), w - Inches(0.7), Inches(0.7),
             head, size=20, bold=True, color=SLATE_800)
        text(s, cx + Inches(0.35), y + Inches(1.7), w - Inches(0.7), Inches(0.9),
             body, size=13, color=SLATE_500, line_spacing=1.55)
        if i < 2:
            arrow_x = cx + w + Inches(0.05)
            text(s, arrow_x, y + Inches(1.05), gap, Inches(0.5),
                 "\u2192", size=32, color=SERENE_300, align=PP_ALIGN.CENTER)

    # Tech stack pill bar
    stack = "Vanilla HTML  \u00b7  Vanilla CSS  \u00b7  Vanilla JS  \u00b7  Google Gemini API"
    rrect(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.55),
          SLATE_50, line_color=SLATE_100, radius=0.3)
    text(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4),
         stack, size=12, bold=True, color=SLATE_700, align=PP_ALIGN.CENTER)
    brand_footer(s)
    page_number(s, n, total)


@slide
def crisis_page(s, n, total):
    bg(s, WHITE)
    # amber-toned banner
    rrect(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(5.8),
          AMBER_50, line_color=AMBER_200, radius=0.04)

    rrect(s, Inches(1.3), Inches(1.25), Inches(0.7), Inches(0.7), AMBER_600, radius=0.25)
    text(s, Inches(1.3), Inches(1.3), Inches(0.7), Inches(0.6),
         "!", size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    eyebrow(s, Inches(2.2), Inches(1.3), "If you need help now", color=AMBER_600)
    title(s, Inches(2.2), Inches(1.55), Inches(9.5),
          "You don\u2019t have to carry this alone.", color=AMBER_800, size=36)

    # Contact tiles
    contacts = [
        ("CALL OR TEXT",  "988",              "Suicide & Crisis Lifeline \u00b7 24/7 \u00b7 US"),
        ("TEXT HOME TO",  "741741",           "Crisis Text Line \u00b7 24/7 \u00b7 Free"),
        ("CALL",          "1-800-662-HELP",   "SAMHSA National Helpline"),
        ("CALL",          "1-866-488-7386",   "The Trevor Project \u00b7 LGBTQ+ youth"),
    ]
    x = Inches(1.3)
    y = Inches(3.05)
    w = Inches(5.35)
    h = Inches(1.6)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    for i, (label, number, desc) in enumerate(contacts):
        r = i // 2
        c = i % 2
        cx = x + (w + gap_x) * c
        cy = y + (h + gap_y) * r
        rrect(s, cx, cy, w, h, WHITE, line_color=AMBER_200, radius=0.08)
        text(s, cx + Inches(0.35), cy + Inches(0.25), w - Inches(0.7), Inches(0.3),
             label, size=10, bold=True, color=AMBER_800)
        text(s, cx + Inches(0.35), cy + Inches(0.55), w - Inches(0.7), Inches(0.7),
             number, size=26, bold=True, color=SLATE_800, line_spacing=1.1)
        text(s, cx + Inches(0.35), cy + Inches(1.15), w - Inches(0.7), Inches(0.4),
             desc, size=11, color=SLATE_500)
    # Footer note
    text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.3),
         "Outside the US? findahelpline.com lists verified crisis services in 130+ countries.",
         size=11, color=SLATE_500, align=PP_ALIGN.CENTER)
    page_number(s, n, total)


@slide
def limits(s, n, total):
    bg(s, WHITE)
    eyebrow(s, Inches(0.9), Inches(0.85), "Honest limits")
    title(s, Inches(0.9), Inches(1.15), Inches(11),
          "What haven is not.")
    subtitle(s, Inches(0.9), Inches(2.2), Inches(11),
             "A tool that refuses to name its limits isn\u2019t safe. So here they are, plainly.")

    # Two columns: NOT / IS
    col_y = Inches(3.25)
    col_h = Inches(3.5)
    # NOT column
    rrect(s, Inches(0.9), col_y, Inches(5.8), col_h, RGBColor(0xfe, 0xf2, 0xf2),
          line_color=RGBColor(0xfe, 0xcd, 0xd3), radius=0.06)
    text(s, Inches(1.2), col_y + Inches(0.3), Inches(5.5), Inches(0.5),
         "haven is NOT", size=16, bold=True, color=RGBColor(0xb9, 0x1c, 0x1c))
    not_items = [
        "A replacement for licensed therapy",
        "A diagnostic tool \u2014 it cannot diagnose",
        "A prescriber \u2014 no medical advice",
        "A crisis intervention service",
        "A friend who remembers you next session",
    ]
    for i, item in enumerate(not_items):
        text(s, Inches(1.2), col_y + Inches(1.0) + Inches(0.45) * i, Inches(5.3), Inches(0.45),
             "\u2717   " + item, size=13, color=SLATE_700, line_spacing=1.4)

    # IS column
    rrect(s, Inches(7.1), col_y, Inches(5.3), col_h, SERENE_50,
          line_color=SERENE_200, radius=0.06)
    text(s, Inches(7.4), col_y + Inches(0.3), Inches(4.8), Inches(0.5),
         "haven IS", size=16, bold=True, color=SERENE_700)
    is_items = [
        "A calm place to think out loud",
        "A gentle reflector of what you\u2019re feeling",
        "A source of evidence-based coping tools",
        "A nudge toward professional help when needed",
        "Always-on, private, and free to use",
    ]
    for i, item in enumerate(is_items):
        text(s, Inches(7.4), col_y + Inches(1.0) + Inches(0.45) * i, Inches(4.8), Inches(0.45),
             "\u2713   " + item, size=13, color=SLATE_700, line_spacing=1.4)
    brand_footer(s)
    page_number(s, n, total)


@slide
def thank_you(s, n, total):
    bg(s, SERENE_50)
    oval(s, Inches(-2), Inches(-2), Inches(7), Inches(7), SERENE_100)
    oval(s, Inches(9), Inches(3.5), Inches(7), Inches(7), SERENE_100)

    logo(s, Inches(6.3), Inches(1.25), size=Inches(0.9))
    text(s, Inches(0), Inches(2.3), Inches(13.333), Inches(1.5),
         "Ready to take a breath?", size=56, bold=True, color=SLATE_800,
         align=PP_ALIGN.CENTER, line_spacing=1.1)
    text(s, Inches(0), Inches(3.6), Inches(13.333), Inches(0.8),
         "Open haven \u2014 a gentle space for your mind.",
         size=20, color=SLATE_500, align=PP_ALIGN.CENTER)
    # CTA pill
    cta_w = Inches(2.7)
    cta_x = (prs.slide_width - cta_w) / 2
    rrect(s, cta_x, Inches(4.7), cta_w, Inches(0.7), SERENE_600, radius=0.5)
    text(s, cta_x, Inches(4.83), cta_w, Inches(0.45),
         "Try haven  \u2192", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    text(s, Inches(0), Inches(5.9), Inches(13.333), Inches(0.4),
         "Thank you.", size=22, bold=True, color=SERENE_700,
         align=PP_ALIGN.CENTER)
    text(s, Inches(0), Inches(6.5), Inches(13.333), Inches(0.4),
         "Questions welcome.", size=13, color=SLATE_500, align=PP_ALIGN.CENTER)
    page_number(s, n, total)


# ------------------------------------------------------------------
# Build
# ------------------------------------------------------------------
total = len(SLIDES)
for i, builder in enumerate(SLIDES, start=1):
    s = prs.slides.add_slide(BLANK)
    builder(s, i, total)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "haven.pptx")
prs.save(out)
print(f"Saved {out} with {total} slides.")
